"""Extract content from PowerPoint files (.pptx and legacy .ppt)."""

import os
import struct
from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt


def extract_text_from_shape(shape) -> str:
    """Extract text from a shape, handling text frames."""
    if not shape.has_text_frame:
        return ""
    text_parts = []
    for para in shape.text_frame.paragraphs:
        para_text = ""
        for run in para.runs:
            para_text += run.text
        if para_text.strip():
            text_parts.append(para_text.strip())
    return "\n".join(text_parts)


def extract_table_content(shape) -> List[List[str]]:
    """Extract table as 2D list of strings."""
    table = shape.table
    rows = []
    for row in table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return rows


def _extract_ppt_ole(filepath: str) -> Dict[str, Any]:
    """Extract text from legacy .ppt (OLE compound) files via olefile.

    Older PPT format stores text in PowerPoint Document streams inside
    an OLE compound file. We extract raw text records using olefile.
    """
    import olefile

    ole = olefile.OleFileIO(filepath)
    slides_text: List[str] = []

    # PPT streams: "PowerPoint Document" is the main stream
    # Slides are stored as individual streams named "PPxx" or in the document stream
    try:
        # Try reading the main PowerPoint Document stream
        ppt_stream = ole.openstream('PowerPoint Document')
        raw = ppt_stream.read()
        # Extract UTF-16LE text chunks (PowerPoint stores text as UTF-16)
        slide_idx = 1
        i = 0
        current_slide_texts: List[str] = []

        while i < len(raw) - 1:
            # Look for TextCharsAtom record (type 0x0FA0)
            if i + 8 <= len(raw):
                rec_ver = raw[i] & 0x0F
                rec_type = struct.unpack_from('<H', raw, i + 2)[0] if i + 4 <= len(raw) else 0
                rec_len = struct.unpack_from('<I', raw, i + 4)[0] if i + 8 <= len(raw) else 0

                if rec_type == 0x0FA0 and rec_len > 0 and rec_len < 100000:
                    # TextCharsAtom: extract text
                    try:
                        text_bytes = raw[i + 8:i + 8 + rec_len]
                        # Try UTF-16LE decode
                        text = text_bytes.decode('utf-16-le', errors='ignore')
                        text = text.replace('\r', '\n').replace('\x00', '')
                        if text.strip() and len(text.strip()) > 1:
                            current_slide_texts.append(text.strip())
                    except Exception:
                        pass

                # Move to next record
                i += 8 + rec_len
            else:
                i += 1

            # Heuristic: split slides at record boundaries (0x0FA0 rec_type 0x03F5 = slide persist)
            # Actually, just collect all text; we don't have perfect slide boundaries
            if rec_type == 0x03F5 and current_slide_texts:
                slides_text.append('\n'.join(current_slide_texts))
                slide_idx += 1
                current_slide_texts = []

        # Don't forget last slide's text
        if current_slide_texts:
            slides_text.append('\n'.join(current_slide_texts))

    except Exception:
        slides_text = []

    # Fallback: extract all readable strings from the OLE container
    if not slides_text:
        try:
            all_streams = ole.listdir()
            for stream_path in all_streams:
                try:
                    stream = ole.openstream('/'.join(stream_path))
                    data = stream.read()
                    # Try to decode as UTF-16LE
                    try:
                        text = data.decode('utf-16-le', errors='ignore')
                        text = text.replace('\x00', '')
                        if len(text) > 50:
                            slides_text.append(text[:5000])  # limit per stream
                    except Exception:
                        pass
                except Exception:
                    pass
        except Exception:
            pass

    ole.close()

    return {
        'slides': [{'number': i + 1, 'title': '', 'text': t, 'tables': [], 'image_count': 0, 'notes': ''}
                    for i, t in enumerate(slides_text)],
        'total_slides': len(slides_text) or 1,
        'filename': os.path.basename(filepath),
    }


def extract_pptx(filepath: str) -> Dict[str, Any]:
    """
    Extract all content from a PPTX/PPT file.
    For .pptx: full extraction (text, tables, notes, placeholder structure).
    For .ppt: fallback text extraction via olefile.

    Returns: {
        'slides': [{'number': int, 'title': str, 'text': str, 'tables': [...], 'notes': str}],
        'total_slides': int
    }
    """
    ext = os.path.splitext(filepath)[1].lower()

    if ext == '.ppt':
        return _extract_ppt_ole(filepath)

    # .pptx: standard python-pptx extraction
    prs = Presentation(filepath)
    slides = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_data = {
            'number': slide_idx,
            'title': '',
            'text': '',
            'tables': [],
            'image_count': 0,
            'notes': '',
        }

        texts = []
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE
                slide_data['title'] = extract_text_from_shape(shape)
            elif shape.has_table:
                slide_data['tables'].append(extract_table_content(shape))
            elif shape.shape_type == 13:  # Picture
                slide_data['image_count'] += 1
            else:
                t = extract_text_from_shape(shape)
                if t:
                    texts.append(t)

        slide_data['text'] = '\n'.join(texts)

        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_data['notes'] = slide.notes_slide.notes_text_frame.text.strip()

        slides.append(slide_data)

    return {
        'slides': slides,
        'total_slides': len(slides),
        'filename': os.path.basename(filepath),
    }


def build_text_summary(data: Dict[str, Any]) -> str:
    """Build a readable text summary from extracted PPTX/PPT data."""
    lines = [f"# {data['filename']}", f"共 {data['total_slides']} 张幻灯片\n"]
    for slide in data['slides']:
        lines.append(f"\n## 幻灯片 {slide['number']}")
        if slide.get('title'):
            lines.append(f"### {slide['title']}")
        if slide.get('text'):
            lines.append(slide['text'])
        for table in slide.get('tables', []):
            if table and table[0]:
                lines.append("\n| " + " | ".join(table[0]) + " |")
                lines.append("|" + "|".join(["---"] * len(table[0])) + "|")
                for row in table[1:]:
                    lines.append("| " + " | ".join(row) + " |")
        if slide.get('notes'):
            lines.append(f"\n> 备注: {slide['notes']}")
    return '\n'.join(lines)
